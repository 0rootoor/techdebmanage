from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Request, Form
from fastapi.responses import HTMLResponse, RedirectResponse, StreamingResponse
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware
from pydantic import BaseModel
from sqlalchemy import create_engine, Column, Integer, String, ForeignKey, Date, Boolean, DateTime, text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
from datetime import date, datetime, timedelta
import pandas as pd
import io
import json
import os
import urllib.request
import urllib.error
import smtplib
from email.mime.text import MIMEText
from openpyxl import Workbook as _OpenpyxlWorkbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Font as _XlsxFont, PatternFill as _XlsxPatternFill
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import hashlib
import secrets
import inspect

TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")
templates = Jinja2Templates(directory=TEMPLATES_DIR)
# Jinja2Templates active déjà l'autoescape HTML par défaut pour les fichiers .html.
# On ajoute un filtre tojson pour pouvoir injecter des valeurs en toute sécurité
# dans les attributs onclick (échappées ensuite en HTML via |e).
templates.env.filters["tojson"] = lambda value: json.dumps(value, ensure_ascii=False).replace("</", "<\\/")
# Contournement d'un bug connu d'incompatibilité entre Starlette >= 1.0.0 et Jinja2
# (TypeError: cannot use 'tuple' as a dict key), lié à la clé de cache interne de Jinja2.
# Désactiver le cache de templates supprime la cause du problème, quelle que soit la
# version de Starlette installée — donc plus besoin de dépendre d'un pin de version.
# Coût : le template est relu/recompilé à chaque requête, négligeable pour une appli
# de cette taille. Voir https://github.com/fastapi/fastapi/issues/15197
templates.env.cache = None

# Compatibilité multi-versions de Starlette pour TemplateResponse :
# - Starlette < 0.29  : TemplateResponse(name, context)            (request DANS context)
# - Starlette >= 1.0  : TemplateResponse(request, name, context)   (request en 1er argument, obligatoire)
# On détecte la signature réellement installée à l'exécution pour ne jamais avoir
# à se soucier de la version de Starlette présente dans l'environnement virtuel.
try:
    _TR_PARAMS = list(inspect.signature(templates.TemplateResponse).parameters)
except (TypeError, ValueError):
    _TR_PARAMS = []
_TEMPLATE_RESPONSE_NEW_STYLE = bool(_TR_PARAMS) and _TR_PARAMS[0] == "request"

def render_template(request: Request, name: str, context: dict, status_code: int = 200):
    """Rend un template Jinja2 via Starlette, quelle que soit la version installée."""
    if _TEMPLATE_RESPONSE_NEW_STYLE:
        return templates.TemplateResponse(request, name, context, status_code=status_code)
    try:
        return templates.TemplateResponse(name, {**context, "request": request}, status_code=status_code)
    except TypeError:
        # Filet de sécurité si la détection de signature s'est trompée
        return templates.TemplateResponse(request, name, context, status_code=status_code)

# Configuration de la base de données SQLite
DATABASE_URL = "sqlite:///./tech_debt_v4.db"
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# --- Modèles SQLAlchemy ---

class ProjectModel(Base):
    __tablename__ = "projects"
    id = Column(Integer, primary_key=True, index=True)
    name = Column(String, unique=True, index=True)
    description = Column(String)
    is_pilot = Column(Boolean, default=False, nullable=False)
    app_status = Column(String, default="En projet")
    socle = Column(String, nullable=True)
    framework = Column(String, nullable=True)

    debts = relationship("TechDebtModel", back_populates="project", cascade="all, delete-orphan")

class TechDebtModel(Base):
    __tablename__ = "tech_debts"
    id = Column(Integer, primary_key=True, index=True)
    title = Column(String, index=True)
    category = Column(String, default="Code")
    impact = Column(String, default="Moyen")
    cost_days = Column(Integer)
    status = Column(String, default="Ouverte")
    created_at = Column(Date, default=date.today)
    start_date = Column(Date, nullable=True)
    target_date = Column(Date, nullable=True)
    assignee = Column(String, nullable=True)
    tags = Column(String, nullable=True)  # tags libres séparés par des virgules, ex: "urgent,q3-2026"
    
    project_id = Column(Integer, ForeignKey("projects.id"))
    project = relationship("ProjectModel", back_populates="debts")
    comments = relationship("CommentModel", back_populates="debt", cascade="all, delete-orphan", order_by="CommentModel.created_at")
    links = relationship("DebtLinkModel", back_populates="debt", cascade="all, delete-orphan")

class AuditLogModel(Base):
    __tablename__ = "audit_log"
    id = Column(Integer, primary_key=True, index=True)
    timestamp = Column(DateTime, default=datetime.utcnow)
    username = Column(String)
    entity_type = Column(String)   # "Application" ou "Dette"
    entity_name = Column(String)
    action = Column(String)        # "Création", "Modification", "Suppression", "Changement de statut"
    details = Column(String, nullable=True)

class UserModel(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True)
    password_hash = Column(String)
    role = Column(String, default="contributeur")  # "admin", "contributeur", "lecture_seule"
    created_at = Column(DateTime, default=datetime.utcnow)

class CommentModel(Base):
    __tablename__ = "comments"
    id = Column(Integer, primary_key=True, index=True)
    debt_id = Column(Integer, ForeignKey("tech_debts.id"))
    username = Column(String)
    content = Column(String)
    created_at = Column(DateTime, default=datetime.utcnow)

    debt = relationship("TechDebtModel", back_populates="comments")

class DebtLinkModel(Base):
    __tablename__ = "debt_links"
    id = Column(Integer, primary_key=True, index=True)
    debt_id = Column(Integer, ForeignKey("tech_debts.id"))
    label = Column(String)   # ex: "Jira TECH-123", "PR #456"
    url = Column(String)

    debt = relationship("TechDebtModel", back_populates="links")

class MilestoneModel(Base):
    __tablename__ = "milestones"
    id = Column(Integer, primary_key=True, index=True)
    label = Column(String)
    milestone_date = Column(Date)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=True)  # None = jalon global
    created_by = Column(String, nullable=True)

    project = relationship("ProjectModel")

Base.metadata.create_all(bind=engine)

# Petite migration de compatibilité : si la base existait déjà (avant l'ajout
# de is_pilot), create_all() ne modifie pas les tables existantes. On ajoute
# la colonne manuellement si besoin. À remplacer par Alembic pour une vraie gestion de schéma.
with engine.connect() as conn:
    existing_columns = [row[1] for row in conn.execute(text("PRAGMA table_info(projects)"))]
    if "is_pilot" not in existing_columns:
        conn.execute(text("ALTER TABLE projects ADD COLUMN is_pilot BOOLEAN NOT NULL DEFAULT 0"))
        conn.commit()
    if "app_status" not in existing_columns:
        conn.execute(text("ALTER TABLE projects ADD COLUMN app_status VARCHAR DEFAULT 'En projet'"))
        conn.commit()
    if "socle" not in existing_columns:
        conn.execute(text("ALTER TABLE projects ADD COLUMN socle VARCHAR"))
        conn.commit()
    if "framework" not in existing_columns:
        conn.execute(text("ALTER TABLE projects ADD COLUMN framework VARCHAR"))
        conn.commit()

    existing_debt_columns = [row[1] for row in conn.execute(text("PRAGMA table_info(tech_debts)"))]
    if "start_date" not in existing_debt_columns:
        conn.execute(text("ALTER TABLE tech_debts ADD COLUMN start_date DATE"))
        conn.commit()
    if "tags" not in existing_debt_columns:
        conn.execute(text("ALTER TABLE tech_debts ADD COLUMN tags VARCHAR"))
        conn.commit()

# --- Application FastAPI ---

app = FastAPI(title="Gestion Avancée de la Dette Technique")

# --- Authentification ---
# Comptes individuels avec rôles (admin / contributeur / lecture_seule), mots de passe
# hachés (PBKDF2-HMAC-SHA256, sans dépendance externe) + session signée (cookie).
# Pour un déploiement au-delà d'une petite équipe interne, il faudra remplacer ceci
# par une vraie IAM d'entreprise (SSO/LDAP/OAuth).
SESSION_SECRET_KEY = os.environ.get("TECHDEBT_SECRET_KEY", "dev-secret-key-change-in-production")
DEFAULT_ADMIN_PASSWORD = os.environ.get("TECHDEBT_ADMIN_PASSWORD", "changeme123")
if SESSION_SECRET_KEY == "dev-secret-key-change-in-production":
    print("⚠️  ATTENTION : clé de session par défaut utilisée. Définis la variable d'environnement "
          "TECHDEBT_SECRET_KEY avant tout déploiement au-delà de ton poste.")

ROLES = ["admin", "contributeur", "lecture_seule"]
ROLE_LABELS = {"admin": "Administrateur", "contributeur": "Contributeur", "lecture_seule": "Lecture seule"}

def hash_password(password: str, salt: str = None) -> str:
    salt = salt or secrets.token_hex(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt.encode("utf-8"), 200_000)
    return f"{salt}${digest.hex()}"

def verify_password(password: str, stored_hash: str) -> bool:
    try:
        salt, _ = stored_hash.split("$", 1)
    except ValueError:
        return False
    return secrets.compare_digest(hash_password(password, salt), stored_hash)

app.add_middleware(SessionMiddleware, secret_key=SESSION_SECRET_KEY, session_cookie="techdebt_session")

# Bootstrap : crée un compte admin par défaut si aucun utilisateur n'existe encore.
with SessionLocal() as _bootstrap_db:
    if _bootstrap_db.query(UserModel).count() == 0:
        _bootstrap_db.add(UserModel(username="admin", password_hash=hash_password(DEFAULT_ADMIN_PASSWORD), role="admin"))
        _bootstrap_db.commit()
        if DEFAULT_ADMIN_PASSWORD == "changeme123":
            print("⚠️  ATTENTION : compte admin créé avec le mot de passe par défaut 'changeme123' "
                  "(identifiant : admin). Change-le dès la première connexion, ou définis "
                  "TECHDEBT_ADMIN_PASSWORD avant le premier démarrage.")

# --- Alertes Slack (optionnel) ---
# Si la variable d'environnement TECHDEBT_SLACK_WEBHOOK_URL est définie, un bouton
# permet d'envoyer un résumé des alertes sur Slack. Sans elle, les alertes restent
# visibles uniquement dans l'onglet "Alertes" de l'application.
SLACK_WEBHOOK_URL = os.environ.get("TECHDEBT_SLACK_WEBHOOK_URL", "")

def send_slack_message(text: str) -> bool:
    if not SLACK_WEBHOOK_URL:
        return False
    try:
        data = json.dumps({"text": text}).encode("utf-8")
        req = urllib.request.Request(SLACK_WEBHOOK_URL, data=data, headers={"Content-Type": "application/json"})
        urllib.request.urlopen(req, timeout=5)
        return True
    except Exception as e:
        print(f"Échec de l'envoi Slack : {e}")
        return False

# --- Alertes Microsoft Teams (optionnel) ---
# Nécessite un connecteur "Webhook entrant" configuré sur un canal Teams.
# Variable d'environnement : TECHDEBT_TEAMS_WEBHOOK_URL
TEAMS_WEBHOOK_URL = os.environ.get("TECHDEBT_TEAMS_WEBHOOK_URL", "")

def send_teams_message(title: str, text: str) -> bool:
    if not TEAMS_WEBHOOK_URL:
        return False
    try:
        # Format "MessageCard", toujours pris en charge par les connecteurs de webhook
        # entrant Teams classiques (Office 365 Connectors).
        payload = {
            "@type": "MessageCard",
            "@context": "http://schema.org/extensions",
            "summary": title,
            "themeColor": "0B2545",
            "title": title,
            "text": text.replace("\n", "\n\n"),  # double saut de ligne = nouveau paragraphe en Markdown Teams
        }
        data = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(TEAMS_WEBHOOK_URL, data=data, headers={"Content-Type": "application/json"})
        urllib.request.urlopen(req, timeout=5)
        return True
    except Exception as e:
        print(f"Échec de l'envoi Teams : {e}")
        return False

# --- Alertes par email (optionnel) ---
# Variables d'environnement :
#   TECHDEBT_SMTP_HOST, TECHDEBT_SMTP_PORT (défaut 587), TECHDEBT_SMTP_USER,
#   TECHDEBT_SMTP_PASSWORD, TECHDEBT_SMTP_FROM (défaut = TECHDEBT_SMTP_USER),
#   TECHDEBT_ALERT_EMAILS (destinataires, séparés par des virgules)
SMTP_HOST = os.environ.get("TECHDEBT_SMTP_HOST", "")
SMTP_PORT = int(os.environ.get("TECHDEBT_SMTP_PORT", "587"))
SMTP_USER = os.environ.get("TECHDEBT_SMTP_USER", "")
SMTP_PASSWORD = os.environ.get("TECHDEBT_SMTP_PASSWORD", "")
SMTP_FROM = os.environ.get("TECHDEBT_SMTP_FROM", SMTP_USER)
ALERT_EMAILS = [e.strip() for e in os.environ.get("TECHDEBT_ALERT_EMAILS", "").split(",") if e.strip()]
EMAIL_ALERTS_ENABLED = bool(SMTP_HOST and SMTP_USER and SMTP_PASSWORD and ALERT_EMAILS)

def send_alert_email(subject: str, body: str) -> bool:
    if not EMAIL_ALERTS_ENABLED:
        return False
    try:
        msg = MIMEText(body, "plain", "utf-8")
        msg["Subject"] = subject
        msg["From"] = SMTP_FROM
        msg["To"] = ", ".join(ALERT_EMAILS)
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=10) as server:
            server.starttls()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_FROM, ALERT_EMAILS, msg.as_string())
        return True
    except Exception as e:
        print(f"Échec de l'envoi email : {e}")
        return False

# --- Résumés IA (optionnel) ---
# Si la variable d'environnement TECHDEBT_LLM_API_KEY est définie, un bouton permet
# de générer un résumé en langage naturel des onglets Alertes et Portefeuille, à
# partir des données déjà calculées côté serveur.
#
# Utilise le format d'API "chat completions" compatible OpenAI, qui fonctionne avec
# la plupart des fournisseurs sans SDK dédié : Mistral (Devstral, Mistral Large...),
# Ollama, vLLM, LM Studio, OpenRouter, etc. Aucune dépendance supplémentaire requise
# (simple appel HTTP via urllib, déjà utilisé pour Slack).
#
# Variables d'environnement :
#   TECHDEBT_LLM_API_KEY   : ta clé/token d'accès (obligatoire pour activer la fonctionnalité)
#   TECHDEBT_LLM_BASE_URL  : URL de l'endpoint "chat completions"
#                            (défaut : https://api.mistral.ai/v1/chat/completions)
#   TECHDEBT_LLM_MODEL     : nom du modèle à utiliser (défaut : "devstral-2-latest" —
#                            vérifie le nom exact disponible sur console.mistral.ai/models
#                            si ça ne fonctionne pas, il change de temps en temps)
LLM_API_KEY = os.environ.get("TECHDEBT_LLM_API_KEY", "")
LLM_BASE_URL = os.environ.get("TECHDEBT_LLM_BASE_URL", "https://api.mistral.ai/v1/chat/completions")
LLM_MODEL = os.environ.get("TECHDEBT_LLM_MODEL", "devstral-2-latest")
AI_SUMMARY_ENABLED = bool(LLM_API_KEY)

def generate_ai_summary(digest: str) -> str:
    """Envoie un texte factuel déjà construit côté serveur au modèle configuré, qui le
    reformule en synthèse en langage naturel. Le modèle n'a accès à aucune autre donnée
    que ce texte : il ne peut pas inventer de chiffres qui n'y figurent pas (mais peut
    se tromper en le reformulant — à relire avant usage officiel)."""
    payload = {
        "model": LLM_MODEL,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Tu rédiges des synthèses factuelles et concises pour un comité de pilotage de "
                    "gestion de dette technique. Utilise UNIQUEMENT les données fournies dans le message "
                    "de l'utilisateur, n'invente jamais de chiffres, de noms d'applications ou de dettes. "
                    "Réponds en français, en 3 à 5 phrases, ton professionnel et direct, sans formule "
                    "d'introduction ni de conclusion générique."
                ),
            },
            {"role": "user", "content": digest},
        ],
        "max_tokens": 400,
        "temperature": 0.3,
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        LLM_BASE_URL,
        data=data,
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {LLM_API_KEY}"},
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        result = json.loads(resp.read().decode("utf-8"))
    return result["choices"][0]["message"]["content"].strip()

def require_api_auth(request: Request) -> str:
    """Dépendance pour les endpoints API en lecture ou pour tout utilisateur connecté : lève une 401 JSON si non connecté."""
    if not request.session.get("authenticated"):
        raise HTTPException(status_code=401, detail="Authentification requise")
    return request.session.get("username", "Utilisateur")

def require_contributor(request: Request) -> str:
    """Dépendance pour les actions d'écriture (créer/modifier/supprimer dettes, commentaires, liens...).
    Les comptes en lecture seule sont bloqués."""
    username = require_api_auth(request)
    if request.session.get("role") == "lecture_seule":
        raise HTTPException(status_code=403, detail="Ton compte est en lecture seule : cette action n'est pas autorisée.")
    return username

def require_admin(request: Request) -> str:
    """Dépendance pour les actions d'administration (suppression d'application, gestion des utilisateurs)."""
    username = require_api_auth(request)
    if request.session.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Cette action est réservée aux administrateurs.")
    return username

def log_action(db: Session, username: str, entity_type: str, entity_name: str, action: str, details: str = None):
    """Enregistre une entrée dans l'historique. Ne fait pas de commit : à inclure dans la même transaction que l'action elle-même."""
    entry = AuditLogModel(username=username, entity_type=entity_type, entity_name=entity_name, action=action, details=details)
    db.add(entry)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

@app.get("/login", response_class=HTMLResponse)
def login_page(request: Request):
    if request.session.get("authenticated"):
        return RedirectResponse(url="/", status_code=303)
    return render_template(request, "login.html", {"error": None})

@app.post("/login")
def login_submit(request: Request, username: str = Form(...), password: str = Form(...), db: Session = Depends(get_db)):
    user = db.query(UserModel).filter(UserModel.username == username.strip()).first()
    if user and verify_password(password, user.password_hash):
        request.session["authenticated"] = True
        request.session["username"] = user.username
        request.session["role"] = user.role
        request.session["user_id"] = user.id
        return RedirectResponse(url="/", status_code=303)
    return render_template(
        request,
        "login.html",
        {"error": "Identifiant ou mot de passe incorrect."},
        status_code=401,
    )

@app.get("/logout")
def logout(request: Request):
    request.session.clear()
    return RedirectResponse(url="/login", status_code=303)

# --- API Endpoints : Utilisateurs (admin uniquement) ---

@app.post("/api/users")
def create_user_endpoint(
    username: str,
    password: str,
    role: str = "contributeur",
    db: Session = Depends(get_db),
    admin_user: str = Depends(require_admin),
):
    username = username.strip()
    if not username or not password:
        raise HTTPException(status_code=400, detail="Identifiant et mot de passe requis")
    if role not in ROLES:
        raise HTTPException(status_code=400, detail="Rôle invalide")
    if db.query(UserModel).filter(UserModel.username == username).first():
        raise HTTPException(status_code=400, detail="Cet identifiant existe déjà")
    db_user = UserModel(username=username, password_hash=hash_password(password), role=role)
    db.add(db_user)
    log_action(db, admin_user, "Utilisateur", username, "Création", f"Rôle : {ROLE_LABELS.get(role, role)}")
    db.commit()
    return {"message": "Utilisateur créé avec succès"}

@app.put("/api/users/{user_id}")
def update_user_endpoint(
    user_id: int,
    role: str = None,
    password: str = None,
    db: Session = Depends(get_db),
    admin_user: str = Depends(require_admin),
):
    db_user = db.query(UserModel).filter(UserModel.id == user_id).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
    if role is not None:
        if role not in ROLES:
            raise HTTPException(status_code=400, detail="Rôle invalide")
        if db_user.role == "admin" and role != "admin":
            remaining_admins = db.query(UserModel).filter(UserModel.role == "admin", UserModel.id != user_id).count()
            if remaining_admins == 0:
                raise HTTPException(status_code=400, detail="Impossible de rétrograder le dernier administrateur")
        db_user.role = role
    if password:
        db_user.password_hash = hash_password(password)
    log_action(db, admin_user, "Utilisateur", db_user.username, "Modification")
    db.commit()
    return {"message": "Utilisateur mis à jour"}

@app.delete("/api/users/{user_id}")
def delete_user_endpoint(user_id: int, request: Request, db: Session = Depends(get_db), admin_user: str = Depends(require_admin)):
    db_user = db.query(UserModel).filter(UserModel.id == user_id).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
    if db_user.id == request.session.get("user_id"):
        raise HTTPException(status_code=400, detail="Tu ne peux pas supprimer ton propre compte")
    if db_user.role == "admin":
        remaining_admins = db.query(UserModel).filter(UserModel.role == "admin", UserModel.id != user_id).count()
        if remaining_admins == 0:
            raise HTTPException(status_code=400, detail="Impossible de supprimer le dernier administrateur")
    name = db_user.username
    db.delete(db_user)
    log_action(db, admin_user, "Utilisateur", name, "Suppression")
    db.commit()
    return {"message": "Utilisateur supprimé"}

# --- API Endpoints : Projets ---

@app.post("/api/projects")
def create_project_endpoint(
    name: str,
    description: str = "",
    is_pilot: bool = False,
    app_status: str = "En projet",
    socle: str = "",
    framework: str = "",
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    if not name.strip():
        raise HTTPException(status_code=400, detail="Le nom du projet est requis")
    existing = db.query(ProjectModel).filter(ProjectModel.name == name.strip()).first()
    if existing:
        raise HTTPException(status_code=400, detail="Ce projet existe déjà")
    
    db_project = ProjectModel(
        name=name.strip(),
        description=description.strip(),
        is_pilot=is_pilot,
        app_status=app_status,
        socle=socle.strip() or None,
        framework=framework.strip() or None,
    )
    db.add(db_project)
    log_action(db, user, "Application", db_project.name, "Création",
               f"Statut : {app_status}" + (f", pilote" if is_pilot else ""))
    db.commit()
    return {"message": "Projet créé avec succès"}

@app.put("/api/projects/{project_id}")
def update_project_endpoint(
    project_id: int,
    name: str,
    description: str = "",
    is_pilot: bool = False,
    app_status: str = "En projet",
    socle: str = "",
    framework: str = "",
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    db_project = db.query(ProjectModel).filter(ProjectModel.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Application non trouvée")
    if not name.strip():
        raise HTTPException(status_code=400, detail="Le nom du projet est requis")
    duplicate = db.query(ProjectModel).filter(ProjectModel.name == name.strip(), ProjectModel.id != project_id).first()
    if duplicate:
        raise HTTPException(status_code=400, detail="Une autre application porte déjà ce nom")

    db_project.name = name.strip()
    db_project.description = description.strip()
    db_project.is_pilot = is_pilot
    db_project.app_status = app_status
    db_project.socle = socle.strip() or None
    db_project.framework = framework.strip() or None
    log_action(db, user, "Application", db_project.name, "Modification",
               f"Statut : {app_status}" + (f", pilote" if is_pilot else ""))
    db.commit()
    return {"message": "Application mise à jour avec succès"}

@app.delete("/api/projects/{project_id}")
def delete_project_endpoint(project_id: int, db: Session = Depends(get_db), user: str = Depends(require_admin)):
    db_project = db.query(ProjectModel).filter(ProjectModel.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Application non trouvée")
    name = db_project.name
    debt_count = db.query(TechDebtModel).filter(TechDebtModel.project_id == project_id).count()
    db.delete(db_project)  # cascade="all, delete-orphan" supprime aussi les dettes associées
    log_action(db, user, "Application", name, "Suppression",
               f"{debt_count} dette(s) associée(s) supprimée(s) en cascade" if debt_count else None)
    db.commit()
    return {"message": f"Application supprimée" + (f" ainsi que {debt_count} dette(s) associée(s)" if debt_count else "")}

def _clean_str(value):
    """Nettoie une valeur pandas (NaN, float, etc.) en chaîne, ou None si vide."""
    if pd.isna(value):
        return None
    s = str(value).strip()
    if not s or s.lower() == 'nan':
        return None
    return s

def _parse_excel_date(value):
    """Convertit une date issue d'Excel/CSV (Timestamp, string, etc.) en date Python, ou None."""
    if pd.isna(value):
        return None
    if isinstance(value, (pd.Timestamp, datetime)):
        return value.date() if isinstance(value, datetime) else value.to_pydatetime().date()
    s = str(value).strip()
    if not s or s.lower() == 'nan':
        return None
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%d-%m-%Y"):
        try:
            return datetime.strptime(s, fmt).date()
        except ValueError:
            continue
    return None

VALID_APP_STATUSES = {"En projet", "En développement", "En production", "En maintenance", "Décommissionnée"}
VALID_CATEGORIES = {"Code", "Architecture", "Sécurité", "Documentation", "Tests"}
VALID_IMPACTS = {"Faible", "Moyen", "Élevé"}
VALID_DEBT_STATUSES = {"Ouverte", "En cours", "Résolue"}

@app.post("/api/projects/import")
async def import_projects(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    contents = await file.read()
    try:
        if file.filename.endswith('.csv'):
            df = pd.read_csv(io.BytesIO(contents))
        elif file.filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(io.BytesIO(contents))
        else:
            raise HTTPException(status_code=400, detail="Format non supporté (.csv ou .xlsx requis)")

        if 'name' not in df.columns:
            raise HTTPException(status_code=400, detail="Le fichier doit contenir une colonne 'name'")

        imported_projects = 0
        imported_debts = 0

        for _, row in df.iterrows():
            name = _clean_str(row.get('name'))
            if not name:
                continue

            # --- Champs de l'application ---
            description = _clean_str(row.get('description')) or ''
            is_pilot_raw = (_clean_str(row.get('is_pilot')) or '').lower()
            is_pilot = is_pilot_raw in ('1', 'true', 'oui', 'yes', 'x')
            app_status = _clean_str(row.get('app_status'))
            if app_status not in VALID_APP_STATUSES:
                app_status = "En projet"
            socle = _clean_str(row.get('socle'))
            framework = _clean_str(row.get('framework'))

            db_project = db.query(ProjectModel).filter(ProjectModel.name == name).first()
            if not db_project:
                db_project = ProjectModel(
                    name=name, description=description, is_pilot=is_pilot,
                    app_status=app_status, socle=socle, framework=framework,
                )
                db.add(db_project)
                db.flush()  # pour obtenir db_project.id avant le commit final
                imported_projects += 1

            # --- Champ(s) de dette technique (optionnels, sur la même ligne) ---
            debt_title = _clean_str(row.get('debt_title'))
            if debt_title:
                category = _clean_str(row.get('debt_category'))
                if category not in VALID_CATEGORIES:
                    category = "Code"
                impact = _clean_str(row.get('debt_impact'))
                if impact not in VALID_IMPACTS:
                    impact = "Moyen"
                status = _clean_str(row.get('debt_status'))
                if status not in VALID_DEBT_STATUSES:
                    status = "Ouverte"
                cost_days_raw = _clean_str(row.get('debt_cost_days'))
                try:
                    cost_days = int(float(cost_days_raw)) if cost_days_raw else 1
                except ValueError:
                    cost_days = 1
                assignee = _clean_str(row.get('debt_assignee'))
                start_date = _parse_excel_date(row.get('debt_start_date'))
                target_date = _parse_excel_date(row.get('debt_target_date'))
                tags = normalize_tags(_clean_str(row.get('debt_tags')) or "")

                db_debt = TechDebtModel(
                    title=debt_title, category=category, impact=impact, status=status,
                    cost_days=cost_days, assignee=assignee,
                    start_date=start_date, target_date=target_date, tags=tags or None,
                    project=db_project,
                )
                db.add(db_debt)
                imported_debts += 1

        db.commit()
        message = f"{imported_projects} application(s) importée(s)"
        if imported_debts:
            message += f" et {imported_debts} dette(s) importée(s)"
        message += " avec succès !"
        if imported_projects or imported_debts:
            log_action(db, user, "Import", file.filename, "Import fichier",
                       f"{imported_projects} application(s), {imported_debts} dette(s)")
            db.commit()
        return {"message": message}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Erreur : {str(e)}")

# --- API Endpoints : Dettes Techniques (CRUD complet) ---

def normalize_tags(raw_tags: str) -> str:
    """Nettoie une liste de tags séparés par des virgules : espaces retirés, vides ignorés,
    doublons supprimés (insensible à la casse), ordre d'origine conservé."""
    if not raw_tags:
        return ""
    seen = set()
    cleaned = []
    for tag in raw_tags.split(","):
        tag = tag.strip()
        if tag and tag.lower() not in seen:
            seen.add(tag.lower())
            cleaned.append(tag)
    return ",".join(cleaned)

@app.post("/api/debts")
def create_debt_endpoint(
    project_id: int,
    title: str,
    category: str,
    impact: str,
    cost_days: int,
    assignee: str = "",
    start_date: str = "",
    target_date: str = "",
    tags: str = "",
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    start = datetime.strptime(start_date, "%Y-%m-%d").date() if start_date else None
    target = datetime.strptime(target_date, "%Y-%m-%d").date() if target_date else None
    db_debt = TechDebtModel(
        title=title,
        category=category,
        impact=impact,
        cost_days=cost_days,
        assignee=assignee if assignee else None,
        start_date=start,
        target_date=target,
        tags=normalize_tags(tags) or None,
        project_id=project_id
    )
    db.add(db_debt)
    log_action(db, user, "Dette", title, "Création", f"{category} / {impact} / {cost_days}j")
    db.commit()
    return {"message": "Dette ajoutée avec succès"}

@app.put("/api/debts/{debt_id}")
def update_debt_endpoint(
    debt_id: int,
    project_id: int,
    title: str,
    category: str,
    impact: str,
    cost_days: int,
    assignee: str = "",
    start_date: str = "",
    target_date: str = "",
    tags: str = "",
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    db_debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not db_debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    
    start = datetime.strptime(start_date, "%Y-%m-%d").date() if start_date else None
    target = datetime.strptime(target_date, "%Y-%m-%d").date() if target_date else None
    db_debt.project_id = project_id
    db_debt.title = title
    db_debt.category = category
    db_debt.impact = impact
    db_debt.cost_days = cost_days
    db_debt.assignee = assignee if assignee else None
    db_debt.start_date = start
    db_debt.target_date = target
    db_debt.tags = normalize_tags(tags) or None
    log_action(db, user, "Dette", title, "Modification", f"{category} / {impact} / {cost_days}j")
    db.commit()
    return {"message": "Dette mise à jour avec succès"}

@app.delete("/api/debts/{debt_id}")
def delete_debt_endpoint(debt_id: int, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    db_debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not db_debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    title = db_debt.title
    db.delete(db_debt)
    log_action(db, user, "Dette", title, "Suppression")
    db.commit()
    return {"message": "Dette supprimée"}

@app.patch("/api/debts/{debt_id}/status")
def update_debt_status(debt_id: int, status: str, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    db_debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not db_debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    old_status = db_debt.status
    db_debt.status = status
    log_action(db, user, "Dette", db_debt.title, "Changement de statut", f"{old_status} → {status}")
    db.commit()
    return {"message": "Statut mis à jour"}


# --- API Endpoints : Commentaires et liens externes (par dette) ---

@app.get("/api/debts/{debt_id}/comments")
def get_comments(debt_id: int, db: Session = Depends(get_db), user: str = Depends(require_api_auth)):
    debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    return [
        {"id": c.id, "username": c.username, "content": c.content, "created_at": c.created_at.strftime("%Y-%m-%d %H:%M")}
        for c in debt.comments
    ]

@app.post("/api/debts/{debt_id}/comments")
def add_comment(debt_id: int, content: str, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    if not content.strip():
        raise HTTPException(status_code=400, detail="Le commentaire ne peut pas être vide")
    comment = CommentModel(debt_id=debt_id, username=user, content=content.strip())
    db.add(comment)
    log_action(db, user, "Dette", debt.title, "Commentaire", content.strip()[:100])
    db.commit()
    return {"message": "Commentaire ajouté"}

@app.delete("/api/comments/{comment_id}")
def delete_comment(comment_id: int, request: Request, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    comment = db.query(CommentModel).filter(CommentModel.id == comment_id).first()
    if not comment:
        raise HTTPException(status_code=404, detail="Commentaire non trouvé")
    if comment.username != user and request.session.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Tu ne peux supprimer que tes propres commentaires")
    db.delete(comment)
    db.commit()
    return {"message": "Commentaire supprimé"}

@app.get("/api/debts/{debt_id}/links")
def get_links(debt_id: int, db: Session = Depends(get_db), user: str = Depends(require_api_auth)):
    debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    return [{"id": l.id, "label": l.label, "url": l.url} for l in debt.links]

@app.post("/api/debts/{debt_id}/links")
def add_link(debt_id: int, label: str, url: str, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    debt = db.query(TechDebtModel).filter(TechDebtModel.id == debt_id).first()
    if not debt:
        raise HTTPException(status_code=404, detail="Dette non trouvée")
    if not label.strip() or not url.strip():
        raise HTTPException(status_code=400, detail="Le libellé et l'URL sont requis")
    if not (url.strip().startswith("http://") or url.strip().startswith("https://")):
        raise HTTPException(status_code=400, detail="L'URL doit commencer par http:// ou https://")
    link = DebtLinkModel(debt_id=debt_id, label=label.strip(), url=url.strip())
    db.add(link)
    log_action(db, user, "Dette", debt.title, "Lien ajouté", label.strip())
    db.commit()
    return {"message": "Lien ajouté"}

@app.delete("/api/links/{link_id}")
def delete_link(link_id: int, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    link = db.query(DebtLinkModel).filter(DebtLinkModel.id == link_id).first()
    if not link:
        raise HTTPException(status_code=404, detail="Lien non trouvé")
    db.delete(link)
    db.commit()
    return {"message": "Lien supprimé"}


# --- API Endpoints : Jalons (vue Gantt) ---

@app.post("/api/milestones")
def create_milestone(
    label: str,
    milestone_date: str,
    project_id: int = None,
    db: Session = Depends(get_db),
    user: str = Depends(require_contributor),
):
    if not label.strip():
        raise HTTPException(status_code=400, detail="Le libellé du jalon est requis")
    try:
        m_date = datetime.strptime(milestone_date, "%Y-%m-%d").date()
    except ValueError:
        raise HTTPException(status_code=400, detail="Date invalide")
    milestone = MilestoneModel(label=label.strip(), milestone_date=m_date, project_id=project_id, created_by=user)
    db.add(milestone)
    log_action(db, user, "Jalon", label.strip(), "Création", m_date.isoformat())
    db.commit()
    return {"message": "Jalon ajouté"}

@app.delete("/api/milestones/{milestone_id}")
def delete_milestone(milestone_id: int, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    milestone = db.query(MilestoneModel).filter(MilestoneModel.id == milestone_id).first()
    if not milestone:
        raise HTTPException(status_code=404, detail="Jalon non trouvé")
    label = milestone.label
    db.delete(milestone)
    log_action(db, user, "Jalon", label, "Suppression")
    db.commit()
    return {"message": "Jalon supprimé"}


def _build_gantt_export_rows(db: Session):
    """Reconstruit les mêmes lignes que la vue Gantt HTML (dette, dates, impact, statut, pilote)."""
    debts = db.query(TechDebtModel).all()
    rows = []
    for d in debts:
        start = d.start_date or d.created_at or date.today()
        if d.target_date:
            end = d.target_date
        else:
            end = start + timedelta(days=max(d.cost_days, 1))
        if end < start:
            end = start
        rows.append({
            "app": d.project.name if d.project else "Inconnue",
            "title": d.title,
            "start": start,
            "end": end,
            "impact": d.impact,
            "status": d.status,
            "cost_days": d.cost_days,
            "pilot": bool(d.project and d.project.is_pilot),
        })
    rows.sort(key=lambda r: (not r["pilot"], r["start"]))
    return rows


def _export_gantt_xlsx(rows) -> io.BytesIO:
    wb = _OpenpyxlWorkbook()
    ws = wb.active
    ws.title = "Gantt"

    headers = ["Application", "Dette", "Début", "Fin", "Durée (jours)", "Impact", "Statut", "Pilote"]
    ws.append(headers)
    header_fill = _XlsxPatternFill("solid", fgColor="0B2545")
    for cell in ws[1]:
        cell.font = _XlsxFont(bold=True, color="FFFFFF")
        cell.fill = header_fill

    min_date = min((r["start"] for r in rows), default=date.today())

    for r in rows:
        duration = max((r["end"] - r["start"]).days, 1)
        ws.append([
            r["app"], r["title"], r["start"], r["end"], duration, r["impact"], r["status"],
            "Oui" if r["pilot"] else "Non",
        ])

    last_row = 1 + len(rows)
    for col_idx, width in enumerate([22, 30, 12, 12, 14, 10, 12, 8], start=1):
        ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = width
    for row in ws.iter_rows(min_row=2, max_row=last_row, min_col=3, max_col=4):
        for cell in row:
            cell.number_format = "YYYY-MM-DD"

    if rows:
        # Colonne technique cachée : décalage en jours depuis la date la plus ancienne,
        # nécessaire pour construire un "faux" Gantt avec un graphique en barres empilées
        # (technique standard Excel : une 1ère série invisible pour le décalage, une 2e
        # série visible pour la durée).
        ws.cell(row=1, column=10, value="Décalage (jours)")
        for i, r in enumerate(rows, start=2):
            ws.cell(row=i, column=10, value=(r["start"] - min_date).days)
        ws.column_dimensions["J"].hidden = True

        chart = BarChart()
        chart.type = "bar"
        chart.grouping = "stacked"
        chart.overlap = 100
        chart.title = "Vue Gantt — Dette technique"
        chart.height = max(8, min(2 + 0.5 * len(rows), 24))
        chart.width = 30
        chart.y_axis.title = None
        chart.x_axis.title = "Jours depuis le " + min_date.strftime("%Y-%m-%d")

        cats = Reference(ws, min_col=2, min_row=2, max_row=last_row)
        data_offset = Reference(ws, min_col=10, min_row=1, max_row=last_row)
        data_duration = Reference(ws, min_col=5, min_row=1, max_row=last_row)
        chart.add_data(data_offset, titles_from_data=True)
        chart.add_data(data_duration, titles_from_data=True)
        chart.set_categories(cats)

        chart.series[0].graphicalProperties.noFill = True
        chart.series[0].graphicalProperties.line.noFill = True
        chart.series[1].graphicalProperties.solidFill = "1256A3"
        chart.legend = None

        ws.add_chart(chart, f"A{last_row + 3}")

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer


def _export_gantt_pptx(rows) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_title(slide, text):
        box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.4), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0B, 0x25, 0x45)

    if not rows:
        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, "Vue Gantt — Dette technique")
        box = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(10), Inches(1))
        box.text_frame.paragraphs[0].text = "Aucune dette à représenter."
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    min_date = min(r["start"] for r in rows)
    max_date = max(r["end"] for r in rows)
    total_days = max((max_date - min_date).days, 1)

    LABEL_LEFT = Inches(0.4)
    LABEL_WIDTH = Inches(3.0)
    CHART_LEFT = Inches(3.5)
    CHART_WIDTH = Emu(Inches(9.3))
    ROW_HEIGHT = Inches(0.32)
    TOP_START = Inches(1.1)
    ROWS_PER_SLIDE = 16

    px_per_day = CHART_WIDTH / total_days
    impact_colors = {
        "Faible": RGBColor(0x15, 0x7A, 0x5C),
        "Moyen": RGBColor(0xB5, 0x73, 0x0A),
        "Élevé": RGBColor(0xC0, 0x36, 0x2C),
    }

    for batch_start in range(0, len(rows), ROWS_PER_SLIDE):
        batch = rows[batch_start:batch_start + ROWS_PER_SLIDE]
        slide = prs.slides.add_slide(blank_layout)
        suffix = "" if batch_start == 0 else f" (suite {batch_start // ROWS_PER_SLIDE + 1})"
        add_title(slide, f"Vue Gantt — Dette technique{suffix}")

        for i, r in enumerate(batch):
            top = TOP_START + i * ROW_HEIGHT

            label_box = slide.shapes.add_textbox(LABEL_LEFT, top, LABEL_WIDTH, ROW_HEIGHT)
            p = label_box.text_frame.paragraphs[0]
            pilot_marker = "⭐ " if r["pilot"] else ""
            p.text = f"{pilot_marker}{r['title']} ({r['app']})"
            p.font.size = Pt(9)

            offset_days = (r["start"] - min_date).days
            duration_days = max((r["end"] - r["start"]).days, 1)
            bar_left = int(CHART_LEFT) + int(px_per_day * offset_days)
            bar_width = max(int(px_per_day * duration_days), Emu(Inches(0.12)))

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, int(top) + Pt(2), bar_width, int(ROW_HEIGHT) - Pt(4)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = impact_colors.get(r["impact"], RGBColor(0x12, 0x56, 0xA3))
            shape.line.fill.background()
            tf = shape.text_frame
            tf.paragraphs[0].text = r["status"]
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf.word_wrap = False

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


@app.get("/api/gantt/export")
def export_gantt_endpoint(format: str = "xlsx", db: Session = Depends(get_db), user: str = Depends(require_api_auth)):
    rows = _build_gantt_export_rows(db)

    if format == "pptx":
        buffer = _export_gantt_pptx(rows)
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        filename = "gantt_dette_technique.pptx"
    elif format == "xlsx":
        buffer = _export_gantt_xlsx(rows)
        media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        filename = "gantt_dette_technique.xlsx"
    else:
        raise HTTPException(status_code=400, detail="Format invalide (attendu : 'xlsx' ou 'pptx')")

    return StreamingResponse(
        buffer,
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.get("/api/debts/export")
def export_debts(
    ids: str = "",
    format: str = "xlsx",
    db: Session = Depends(get_db),
    user: str = Depends(require_api_auth),
):
    query = db.query(TechDebtModel)
    if ids:
        id_list = [int(x) for x in ids.split(",") if x.strip().isdigit()]
        if id_list:
            query = query.filter(TechDebtModel.id.in_(id_list))
    debts = query.all()

    rows = []
    for d in debts:
        rows.append({
            "Application": d.project.name if d.project else "",
            "Statut app": d.project.app_status if d.project else "",
            "Socle": d.project.socle if d.project and d.project.socle else "",
            "Framework": d.project.framework if d.project and d.project.framework else "",
            "Pilote": "Oui" if d.project and d.project.is_pilot else "Non",
            "Titre dette": d.title,
            "Catégorie": d.category,
            "Impact": d.impact,
            "Statut dette": d.status,
            "Charge (jours)": d.cost_days,
            "Responsable": d.assignee or "",
            "Tags": d.tags or "",
            "Date de début": d.start_date.isoformat() if d.start_date else "",
            "Date cible": d.target_date.isoformat() if d.target_date else "",
        })
    df = pd.DataFrame(rows, columns=[
        "Application", "Statut app", "Socle", "Framework", "Pilote", "Titre dette",
        "Catégorie", "Impact", "Statut dette", "Charge (jours)", "Responsable", "Tags",
        "Date de début", "Date cible",
    ])

    buffer = io.BytesIO()
    if format == "csv":
        df.to_csv(buffer, index=False, encoding="utf-8-sig")
        media_type = "text/csv"
        filename = "registre_dette_technique.csv"
    else:
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Dettes")
            worksheet = writer.sheets["Dettes"]
            for col_idx, col_name in enumerate(df.columns, start=1):
                max_len = max([len(str(col_name))] + [len(str(v)) for v in df[col_name].astype(str)]) if len(df) else len(col_name)
                worksheet.column_dimensions[worksheet.cell(row=1, column=col_idx).column_letter].width = min(max_len + 2, 40)
            for cell in worksheet[1]:
                cell.font = cell.font.copy(bold=True)
        media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        filename = "registre_dette_technique.xlsx"

    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.post("/api/alerts/send")
def send_alerts_endpoint(channel: str = "slack", db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    channel_config = {
        "slack": (bool(SLACK_WEBHOOK_URL), "TECHDEBT_SLACK_WEBHOOK_URL"),
        "teams": (bool(TEAMS_WEBHOOK_URL), "TECHDEBT_TEAMS_WEBHOOK_URL"),
        "email": (EMAIL_ALERTS_ENABLED, "TECHDEBT_SMTP_HOST / TECHDEBT_SMTP_USER / TECHDEBT_SMTP_PASSWORD / TECHDEBT_ALERT_EMAILS"),
    }
    if channel not in channel_config:
        raise HTTPException(status_code=400, detail="Canal invalide (attendu : slack, teams ou email)")
    configured, var_names = channel_config[channel]
    if not configured:
        return {
            "sent": False,
            "message": f"Ce canal n'est pas configuré (variable(s) d'environnement {var_names} absente(s)). "
                       "Les alertes restent visibles dans l'onglet Alertes de l'application."
        }

    debts = db.query(TechDebtModel).all()
    open_debts = [d for d in debts if d.status != "Résolue"]
    overdue = [d for d in open_debts if d.target_date and d.target_date < date.today()]
    soon_threshold = date.today() + timedelta(days=7)
    soon = [d for d in open_debts if d.target_date and date.today() <= d.target_date <= soon_threshold]
    stale_pilot = [
        d for d in open_debts
        if d.project and d.project.is_pilot and d.status == "Ouverte"
        and (date.today() - (d.start_date or d.created_at or date.today())).days > 30
    ]

    if not overdue and not soon and not stale_pilot:
        return {"sent": False, "message": "Aucune alerte à signaler pour le moment."}

    title = f"Alertes dette technique — {date.today().isoformat()}"
    lines = [f"*{title}*"]
    if overdue:
        lines.append(f"\n:red_circle: *{len(overdue)} dette(s) en retard*")
        for d in overdue[:10]:
            lines.append(f"• {d.title} ({d.project.name if d.project else '?'}) — échéance {d.target_date.isoformat()}")
    if soon:
        lines.append(f"\n:large_orange_circle: *{len(soon)} échéance(s) dans les 7 prochains jours*")
        for d in soon[:10]:
            lines.append(f"• {d.title} ({d.project.name if d.project else '?'}) — échéance {d.target_date.isoformat()}")
    if stale_pilot:
        lines.append(f"\n:large_blue_circle: *{len(stale_pilot)} dette(s) pilote(s) ouverte(s) depuis plus de 30 jours*")
        for d in stale_pilot[:10]:
            lines.append(f"• {d.title} ({d.project.name if d.project else '?'})")
    text = "\n".join(lines)

    channel_labels = {"slack": "Slack", "teams": "Teams", "email": "email"}
    if channel == "slack":
        ok = send_slack_message(text)
    elif channel == "teams":
        # Le format Slack (*gras*, :emoji:) n'est pas interprété par Teams : on repart d'un texte simple.
        plain_lines = [f"{len(overdue)} dette(s) en retard, {len(soon)} échéance(s) proche(s), {len(stale_pilot)} dette(s) pilote(s) bloquée(s)."]
        for d in (overdue + soon + stale_pilot)[:20]:
            plain_lines.append(f"- {d.title} ({d.project.name if d.project else '?'})")
        ok = send_teams_message(title, "\n".join(plain_lines))
    else:
        plain_lines = [title, ""]
        plain_lines += [l.replace("*", "").replace(":red_circle:", "🔴").replace(":large_orange_circle:", "🟠").replace(":large_blue_circle:", "🔵") for l in lines[1:]]
        ok = send_alert_email(title, "\n".join(plain_lines))

    if ok:
        log_action(db, user, "Alertes", channel_labels[channel], "Envoi", f"{len(overdue)} retard(s), {len(soon)} proche(s), {len(stale_pilot)} pilote(s) bloquée(s)")
        db.commit()
        return {"sent": True, "message": f"Alertes envoyées par {channel_labels[channel]} avec succès."}
    return {"sent": False, "message": f"Échec de l'envoi par {channel_labels[channel]}. Vérifie la configuration et la connexion réseau."}


@app.post("/api/summary/generate")
def generate_summary_endpoint(scope: str, db: Session = Depends(get_db), user: str = Depends(require_contributor)):
    if not AI_SUMMARY_ENABLED:
        raise HTTPException(
            status_code=400,
            detail="Aucune clé API configurée (variable d'environnement TECHDEBT_LLM_API_KEY absente).",
        )

    if scope == "alerts":
        debts = db.query(TechDebtModel).all()
        open_debts = [d for d in debts if d.status != "Résolue"]
        overdue = [d for d in open_debts if d.target_date and d.target_date < date.today()]
        soon_threshold = date.today() + timedelta(days=7)
        soon = [d for d in open_debts if d.target_date and date.today() <= d.target_date <= soon_threshold]
        stale_pilot = [
            d for d in open_debts
            if d.project and d.project.is_pilot and d.status == "Ouverte"
            and (date.today() - (d.start_date or d.created_at or date.today())).days > 30
        ]

        if not overdue and not soon and not stale_pilot:
            return {"summary": "Aucune alerte active actuellement : pas de dette en retard, pas d'échéance proche, et aucune dette pilote bloquée depuis plus de 30 jours."}

        lines = [f"Date du jour : {date.today().isoformat()}", ""]
        lines.append(f"Dettes en retard ({len(overdue)}) :")
        for d in overdue[:15]:
            lines.append(f"- {d.title} | application : {d.project.name if d.project else 'inconnue'} | échéance dépassée le {d.target_date.isoformat()} | impact {d.impact} | charge {d.cost_days}j")
        lines.append("")
        lines.append(f"Échéances dans les 7 prochains jours ({len(soon)}) :")
        for d in soon[:15]:
            lines.append(f"- {d.title} | application : {d.project.name if d.project else 'inconnue'} | échéance le {d.target_date.isoformat()} | impact {d.impact}")
        lines.append("")
        lines.append(f"Dettes pilotes ouvertes depuis plus de 30 jours sans changement de statut ({len(stale_pilot)}) :")
        for d in stale_pilot[:15]:
            lines.append(f"- {d.title} | application : {d.project.name if d.project else 'inconnue'} | responsable : {d.assignee or 'non assigné'}")
        digest = "\n".join(lines)

    elif scope == "portfolio":
        projects = db.query(ProjectModel).order_by(ProjectModel.is_pilot.desc(), ProjectModel.name).all()
        debts = db.query(TechDebtModel).all()
        rows = []
        for p in projects:
            p_debts = [d for d in debts if d.project_id == p.id]
            overdue_count = sum(1 for d in p_debts if d.target_date and d.target_date < date.today() and d.status != "Résolue")
            rows.append({
                "project": p,
                "debt_count": len(p_debts),
                "total_cost": sum(d.cost_days for d in p_debts),
                "overdue_count": overdue_count,
            })
        if not rows:
            return {"summary": "Aucune application enregistrée pour le moment."}

        total_apps = len(rows)
        total_cost = sum(r["total_cost"] for r in rows)
        total_overdue = sum(r["overdue_count"] for r in rows)
        pilot_apps = sum(1 for r in rows if r["project"].is_pilot)
        top_rows = sorted(rows, key=lambda r: -r["total_cost"])[:8]

        lines = [
            f"Nombre total d'applications : {total_apps}",
            f"Applications en mode pilote : {pilot_apps}",
            f"Charge de dette technique totale (toutes applications) : {total_cost} jours",
            f"Nombre d'applications ayant au moins une dette en retard : {total_overdue}",
            "",
            "Détail des applications avec le plus de charge de dette (jusqu'à 8) :",
        ]
        for r in top_rows:
            p = r["project"]
            lines.append(
                f"- {p.name} | statut : {p.app_status} | socle : {p.socle or 'non renseigné'} | "
                f"framework : {p.framework or 'non renseigné'} | {r['debt_count']} dette(s) | "
                f"{r['total_cost']}j de charge | {r['overdue_count']} en retard"
                + (" | application pilote" if p.is_pilot else "")
            )
        digest = "\n".join(lines)

    else:
        raise HTTPException(status_code=400, detail="Portée invalide (attendu : 'alerts' ou 'portfolio')")

    try:
        summary_text = generate_ai_summary(digest)
    except urllib.error.HTTPError as e:
        try:
            error_body = json.loads(e.read().decode("utf-8"))
            error_detail = error_body.get("message") or error_body.get("error", {}).get("message") or str(error_body)
        except Exception:
            error_detail = str(e)
        raise HTTPException(status_code=502, detail=f"Erreur de l'API IA (HTTP {e.code}) : {error_detail}")
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Erreur lors de l'appel à l'API IA : {e}")

    log_action(db, user, "Résumé IA", scope, "Génération")
    db.commit()
    return {"summary": summary_text}


# --- Interface Frontend Complète ---

@app.get("/", response_class=HTMLResponse)
def read_root(request: Request, db: Session = Depends(get_db)):
    if not request.session.get("authenticated") or "role" not in request.session:
        # "role" absent : session issue de l'ancienne authentification par mot de passe
        # partagé (avant la refonte en comptes individuels) -> on force une reconnexion.
        request.session.clear()
        return RedirectResponse(url="/login", status_code=303)
    current_user = request.session.get("username", "Utilisateur")
    current_role = request.session.get("role", "lecture_seule")

    projects = db.query(ProjectModel).order_by(ProjectModel.is_pilot.desc(), ProjectModel.name).all()
    debts = db.query(TechDebtModel).all()
    total_cost = sum(d.cost_days for d in debts)
    sorted_debts = sorted(debts, key=lambda x: x.target_date if x.target_date else date.max)

    open_debts = [d for d in debts if d.status != "Résolue"]
    overdue_debts = [d for d in open_debts if d.target_date and d.target_date < date.today()]
    pilot_projects = [p for p in projects if p.is_pilot]

    distinct_socles = sorted({p.socle for p in projects if p.socle})
    distinct_frameworks = sorted({p.framework for p in projects if p.framework})

    all_tags = set()
    for d in debts:
        if d.tags:
            all_tags.update(t.strip() for t in d.tags.split(",") if t.strip())
    distinct_tags = sorted(all_tags)

    # Agrégats pour les graphiques (calculés côté serveur, injectés en JSON dans le template)
    categories = ["Code", "Architecture", "Sécurité", "Documentation", "Tests"]
    category_labels, category_counts = [], []
    for cat in categories:
        count = sum(1 for d in debts if d.category == cat)
        if count > 0:
            category_labels.append(cat)
            category_counts.append(count)

    impact_order = ["Faible", "Moyen", "Élevé"]
    impact_counts = [sum(1 for d in debts if d.impact == level) for level in impact_order]

    status_order = ["Ouverte", "En cours", "Résolue"]
    status_counts = [sum(1 for d in debts if d.status == s) for s in status_order]

    cost_by_category = []
    for cat in category_labels:
        cost_by_category.append(sum(d.cost_days for d in debts if d.category == cat))

    # Répartition des applications par socle / framework / statut
    from collections import Counter
    socle_counter = Counter(p.socle for p in projects if p.socle)
    socle_labels = sorted(socle_counter, key=lambda k: -socle_counter[k])
    socle_counts = [socle_counter[k] for k in socle_labels]

    framework_counter = Counter(p.framework for p in projects if p.framework)
    framework_labels = sorted(framework_counter, key=lambda k: -framework_counter[k])
    framework_counts = [framework_counter[k] for k in framework_labels]

    app_status_order = ["En projet", "En développement", "En production", "En maintenance", "Décommissionnée"]
    app_status_labels, app_status_counts = [], []
    for s in app_status_order:
        count = sum(1 for p in projects if p.app_status == s)
        if count > 0:
            app_status_labels.append(s)
            app_status_counts.append(count)

    chart_data = {
        "categoryLabels": category_labels,
        "categoryCounts": category_counts,
        "costByCategory": cost_by_category,
        "impactLabels": impact_order,
        "impactCounts": impact_counts,
        "statusLabels": status_order,
        "statusCounts": status_counts,
        "socleLabels": socle_labels,
        "socleCounts": socle_counts,
        "frameworkLabels": framework_labels,
        "frameworkCounts": framework_counts,
        "appStatusLabels": app_status_labels,
        "appStatusCounts": app_status_counts,
    }

    # Données pour la vue Gantt : une barre par dette, du jour de création
    # jusqu'à la date cible (ou, à défaut, une estimation création + charge en jours).
    gantt_rows = []
    for d in debts:
        start = d.start_date or d.created_at or date.today()
        if d.target_date:
            end = d.target_date
            estimated = False
        else:
            end = start + timedelta(days=max(d.cost_days, 1))
            estimated = True
        if end < start:
            end = start
        gantt_rows.append({
            "id": d.id,
            "title": d.title,
            "project": d.project.name if d.project else "Inconnu",
            "isPilot": bool(d.project and d.project.is_pilot),
            "status": d.status,
            "impact": d.impact,
            "assignee": d.assignee or "Non assigné",
            "costDays": d.cost_days,
            "start": start.isoformat(),
            "end": end.isoformat(),
            "estimated": estimated,
        })
    # Tri : applications pilotes d'abord, puis par date de début
    gantt_rows.sort(key=lambda r: (not r["isPilot"], r["start"]))

    # Jalons de la vue Gantt
    milestones = db.query(MilestoneModel).order_by(MilestoneModel.milestone_date).all()
    milestones_data = [
        {
            "id": m.id,
            "label": m.label,
            "date": m.milestone_date.isoformat(),
            "project": m.project.name if m.project else None,
        }
        for m in milestones
    ]

    # Vue "portefeuille applicatif" : agrégats par application plutôt que par dette
    portfolio_rows = []
    for p in projects:
        p_debts = [d for d in debts if d.project_id == p.id]
        p_overdue = sum(1 for d in p_debts if d.target_date and d.target_date < date.today() and d.status != "Résolue")
        portfolio_rows.append({
            "project": p,
            "debt_count": len(p_debts),
            "total_cost": sum(d.cost_days for d in p_debts),
            "open_count": sum(1 for d in p_debts if d.status == "Ouverte"),
            "in_progress_count": sum(1 for d in p_debts if d.status == "En cours"),
            "resolved_count": sum(1 for d in p_debts if d.status == "Résolue"),
            "overdue_count": p_overdue,
        })
    portfolio_rows.sort(key=lambda r: (not r["project"].is_pilot, -r["total_cost"]))

    # Alertes : échéances dépassées, échéances proches (7 jours), dettes pilotes bloquées (30 jours)
    SOON_DAYS = 7
    STALE_PILOT_DAYS = 30
    soon_threshold = date.today() + timedelta(days=SOON_DAYS)
    alerts_overdue = [d for d in overdue_debts]
    alerts_soon = [
        d for d in open_debts
        if d.target_date and date.today() <= d.target_date <= soon_threshold
    ]
    alerts_stale_pilot = [
        d for d in open_debts
        if d.project and d.project.is_pilot and d.status == "Ouverte"
        and (date.today() - (d.start_date or d.created_at or date.today())).days > STALE_PILOT_DAYS
    ]

    # Historique récent (200 dernières actions, les plus récentes en premier)
    recent_audit_log = db.query(AuditLogModel).order_by(AuditLogModel.timestamp.desc()).limit(200).all()
    all_users = db.query(UserModel).order_by(UserModel.username).all() if current_role == "admin" else []

    return render_template(
        request,
        "index.html",
        {
            "current_user": current_user,
            "current_role": current_role,
            "role_labels": ROLE_LABELS,
            "all_users": all_users,
            "projects": projects,
            "debts": debts,
            "sorted_debts": sorted_debts,
            "total_cost": total_cost,
            "open_debts_count": len(open_debts),
            "overdue_count": len(overdue_debts),
            "pilot_count": len(pilot_projects),
            "distinct_socles": distinct_socles,
            "distinct_frameworks": distinct_frameworks,
            "distinct_tags": distinct_tags,
            "today": date.today(),
            "chart_data_json": json.dumps(chart_data, ensure_ascii=False).replace("</", "<\\/"),
            "gantt_data_json": json.dumps(gantt_rows, ensure_ascii=False).replace("</", "<\\/"),
            "milestones_json": json.dumps(milestones_data, ensure_ascii=False).replace("</", "<\\/"),
            "milestones": milestones,
            "portfolio_rows": portfolio_rows,
            "alerts_overdue": alerts_overdue,
            "alerts_soon": alerts_soon,
            "alerts_stale_pilot": alerts_stale_pilot,
            "soon_days": SOON_DAYS,
            "stale_pilot_days": STALE_PILOT_DAYS,
            "recent_audit_log": recent_audit_log,
            "slack_configured": bool(SLACK_WEBHOOK_URL),
            "teams_configured": bool(TEAMS_WEBHOOK_URL),
            "email_configured": EMAIL_ALERTS_ENABLED,
            "ai_summary_enabled": AI_SUMMARY_ENABLED,
        },
    )